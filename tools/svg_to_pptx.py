#!/usr/bin/env python3
"""
PPT Master - SVG 转 PPTX 工具

将项目中的 SVG 文件批量转换为 PowerPoint 演示文稿。
每个 SVG 文件对应一张幻灯片，SVG 以原生矢量格式嵌入。

用法:
    python3 tools/svg_to_pptx.py <项目路径>
    python3 tools/svg_to_pptx.py <项目路径> -o output.pptx
    python3 tools/svg_to_pptx.py <项目路径> --use-final

示例:
    python3 tools/svg_to_pptx.py examples/ppt169_demo
    python3 tools/svg_to_pptx.py examples/ppt169_demo -o presentation.pptx
    python3 tools/svg_to_pptx.py examples/ppt169_demo --use-final

依赖:
    pip install python-pptx

注意:
    - SVG 以原生矢量格式嵌入，保持可编辑性
    - 需要 PowerPoint 2016+ 才能正确显示
"""

import sys
import os
import argparse
import re
import zipfile
import shutil
import tempfile
from pathlib import Path
from typing import Optional, Tuple, List
from xml.etree import ElementTree as ET

# 检查 python-pptx 是否已安装
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("错误: 缺少 python-pptx 库")
    print("请运行: pip install python-pptx")
    sys.exit(1)

# 导入项目工具模块
sys.path.insert(0, str(Path(__file__).parent))
try:
    from project_utils import get_project_info
    from config import CANVAS_FORMATS
except ImportError:
    CANVAS_FORMATS = {
        'ppt169': {'name': 'PPT 16:9', 'dimensions': '1280×720', 'viewbox': '0 0 1280 720'},
    }

    def get_project_info(path):
        return {'format': 'unknown', 'name': Path(path).name}

# 导入动画模块
try:
    from pptx_animations import create_transition_xml, TRANSITIONS
    ANIMATIONS_AVAILABLE = True
except ImportError:
    ANIMATIONS_AVAILABLE = False
    TRANSITIONS = {}

# SVG 转 PNG 库检测（用于 Office 兼容模式）
# 优先使用 CairoSVG（渲染质量更好），降级到 svglib
PNG_RENDERER = None  # 'cairosvg' | 'svglib' | None

try:
    import cairosvg
    PNG_RENDERER = 'cairosvg'
except ImportError:
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        PNG_RENDERER = 'svglib'
    except ImportError:
        pass

def get_png_renderer_info() -> tuple:
    """获取 PNG 渲染器信息"""
    if PNG_RENDERER == 'cairosvg':
        return ('cairosvg', '(渐变/滤镜完整)', None)
    elif PNG_RENDERER == 'svglib':
        return ('svglib', '(部分渐变可能丢失)', '安装 cairosvg 可获得更好效果: pip install cairosvg')
    else:
        return (None, '(未安装)', '安装方法: pip install cairosvg 或 pip install svglib reportlab')


# EMU 转换常量
EMU_PER_INCH = 914400
EMU_PER_PIXEL = EMU_PER_INCH / 96

# XML 命名空间
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'asvg': 'http://schemas.microsoft.com/office/drawing/2016/SVG/main',
}

# 注册命名空间
for prefix, uri in NAMESPACES.items():
    ET.register_namespace(prefix, uri)


def get_slide_dimensions(canvas_format: str, custom_pixels: Optional[Tuple[int, int]] = None) -> Tuple[int, int]:
    """获取幻灯片尺寸（EMU 单位）"""
    if custom_pixels:
        width_px, height_px = custom_pixels
    else:
        if canvas_format not in CANVAS_FORMATS:
            canvas_format = 'ppt169'
        
        dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
        match = re.match(r'(\d+)[×x](\d+)', dimensions)
        if match:
            width_px = int(match.group(1))
            height_px = int(match.group(2))
        else:
            width_px, height_px = 1280, 720
    
    return int(width_px * EMU_PER_PIXEL), int(height_px * EMU_PER_PIXEL)


def get_pixel_dimensions(canvas_format: str, custom_pixels: Optional[Tuple[int, int]] = None) -> Tuple[int, int]:
    """获取画布像素尺寸"""
    if custom_pixels:
        return custom_pixels
    
    if canvas_format not in CANVAS_FORMATS:
        canvas_format = 'ppt169'
    
    dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
    match = re.match(r'(\d+)[×x](\d+)', dimensions)
    if match:
        return int(match.group(1)), int(match.group(2))
    return 1280, 720


def get_viewbox_dimensions(svg_path: Path) -> Optional[Tuple[int, int]]:
    """从 SVG 的 viewBox 提取像素尺寸（返回整数）"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read(2000)
        
        match = re.search(r'viewBox="([^"]+)"', content)
        if not match:
            return None
        
        parts = re.split(r'[\s,]+', match.group(1).strip())
        if len(parts) < 4:
            return None
        
        width = float(parts[2])
        height = float(parts[3])
        if width <= 0 or height <= 0:
            return None
        
        return int(round(width)), int(round(height))
    except Exception:
        return None


def detect_format_from_svg(svg_path: Path) -> Optional[str]:
    """从 SVG 文件的 viewBox 检测画布格式"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read(2000)
        
        match = re.search(r'viewBox="([^"]+)"', content)
        if match:
            viewbox = match.group(1)
            for fmt_key, fmt_info in CANVAS_FORMATS.items():
                if fmt_info['viewbox'] == viewbox:
                    return fmt_key
    except Exception:
        pass
    return None


def convert_svg_to_png(svg_path: Path, png_path: Path, width: int = None, height: int = None) -> bool:
    """
    将 SVG 转换为 PNG
    
    Args:
        svg_path: SVG 文件路径
        png_path: 输出 PNG 文件路径
        width: 输出宽度（像素）
        height: 输出高度（像素）
    
    Returns:
        是否成功转换
    """
    if PNG_RENDERER is None:
        return False
    
    try:
        if PNG_RENDERER == 'cairosvg':
            # 使用 CairoSVG（渲染质量更好）
            cairosvg.svg2png(
                url=str(svg_path),
                write_to=str(png_path),
                output_width=width,
                output_height=height
            )
            return True
        
        elif PNG_RENDERER == 'svglib':
            # 使用 svglib（轻量级，但渐变支持有限）
            drawing = svg2rlg(str(svg_path))
            if drawing is None:
                print(f"  警告: 无法解析 SVG ({svg_path.name})")
                return False
            
            # 渲染为 PNG
            renderPM.drawToFile(
                drawing,
                str(png_path),
                fmt="PNG",
                configPIL={'quality': 95}
            )
            return True
        
    except Exception as e:
        print(f"  警告: SVG 转 PNG 失败 ({svg_path.name}): {e}")
        return False
    
    return False


def find_svg_files(project_path: Path, source: str = 'output') -> Tuple[List[Path], str]:
    """
    查找项目中的 SVG 文件
    
    Args:
        project_path: 项目目录路径
        source: SVG 来源目录
            - 'output': svg_output（原始版本）
            - 'final': svg_final（后处理完成，推荐）
            - 或任意子目录名称
    
    Returns:
        (SVG 文件列表, 实际使用的目录名)
    """
    # 预定义目录映射
    dir_map = {
        'output': 'svg_output',
        'final': 'svg_final',
    }
    
    # 获取目录名（支持预定义别名或直接指定目录名）
    dir_name = dir_map.get(source, source)
    svg_dir = project_path / dir_name
    
    if not svg_dir.exists():
        print(f"  警告: {dir_name} 目录不存在，尝试 svg_output")
        dir_name = 'svg_output'
        svg_dir = project_path / dir_name
    
    if not svg_dir.exists():
        # 直接在指定目录查找
        if project_path.is_dir():
            svg_dir = project_path
            dir_name = project_path.name
        else:
            return [], ''
    
    return sorted(svg_dir.glob('*.svg')), dir_name


def find_notes_files(project_path: Path, svg_files: List[Path] = None) -> dict:
    """
    查找项目中的备注文件
    
    支持两种匹配模式（支持混合匹配）：
    1. 按文件名匹配（优先）：notes/01_封面.md 对应 01_封面.svg
    2. 按序号匹配（向后兼容）：notes/slide01.md 对应第1个 SVG
    
    Args:
        project_path: 项目目录路径
        svg_files: SVG 文件列表（用于按文件名匹配）
    
    Returns:
        字典，key 为 SVG 文件名（不含扩展名），value 为备注内容
    """
    notes_dir = project_path / 'notes'
    notes = {}
    
    if not notes_dir.exists():
        return notes
    
    svg_stems_mapping = {}
    svg_index_mapping = {}
    if svg_files:
        for i, svg_path in enumerate(svg_files, 1):
            svg_stems_mapping[svg_path.stem] = i
            svg_index_mapping[i] = svg_path.stem

    # 收集所有 notes 文件信息
    for notes_file in notes_dir.glob('*.md'):
        try:
            with open(notes_file, 'r', encoding='utf-8') as f:
                content = f.read().strip()
            if not content:
                continue
            
            stem = notes_file.stem

            # 尝试提取序号（向后兼容 slide01.md 格式）
            match = re.search(r'slide[_]?(\d+)', stem)
            if match:
                index = int(match.group(1))
                mapped_stem = svg_index_mapping.get(index)
                if mapped_stem:
                    notes[mapped_stem] = content

            # 按文件名提取（覆盖向后兼容的格式）
            if stem in svg_stems_mapping:
                notes[stem] = content
        except Exception:
            pass
    
    return notes


def markdown_to_plain_text(md_content: str) -> str:
    """
    将 Markdown 备注转换为纯文本（用于 PPTX 备注）
    
    Args:
        md_content: Markdown 格式的备注内容
    
    Returns:
        纯文本内容
    """
    def strip_inline_bold(text: str) -> str:
        # Remove Markdown bold markers while keeping content
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        return text

    lines = []
    for line in md_content.split('\n'):
        # 跳过标题行（# 开头）
        if line.startswith('#'):
            # 提取标题文本
            text = re.sub(r'^#+\s*', '', line).strip()
            text = strip_inline_bold(text)
            if text:
                lines.append(text)
                lines.append('')  # 空行
        # 处理列表项（- 开头）
        elif line.strip().startswith('- '):
            item_text = line.strip()[2:]
            item_text = strip_inline_bold(item_text)
            lines.append('• ' + item_text)
        # 普通行
        elif line.strip():
            text = strip_inline_bold(line.strip())
            lines.append(text)
        else:
            lines.append('')
    
    # 合并连续空行
    result = []
    prev_empty = False
    for line in lines:
        if line == '':
            if not prev_empty:
                result.append(line)
            prev_empty = True
        else:
            result.append(line)
            prev_empty = False
    
    return '\n'.join(result).strip()


def create_notes_slide_xml(slide_num: int, notes_text: str) -> str:
    """
    创建备注幻灯片 XML
    
    Args:
        slide_num: 幻灯片序号
        notes_text: 备注文本（纯文本格式）
    
    Returns:
        备注幻灯片 XML 字符串
    """
    # 转义 XML 特殊字符
    notes_text = notes_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    
    # 将换行转换为 <a:p> 段落
    paragraphs = []
    for para in notes_text.split('\n'):
        if para.strip():
            paragraphs.append(f'''<a:p>
              <a:r>
                <a:rPr lang="zh-CN" dirty="0"/>
                <a:t>{para}</a:t>
              </a:r>
            </a:p>''')
        else:
            paragraphs.append('<a:p><a:endParaRPr lang="zh-CN" dirty="0"/></a:p>')
    
    paragraphs_xml = '\n            '.join(paragraphs) if paragraphs else '<a:p><a:endParaRPr lang="zh-CN" dirty="0"/></a:p>'
    
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
         xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Slide Image Placeholder 1"/>
          <p:cNvSpPr>
            <a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/>
          </p:cNvSpPr>
          <p:nvPr>
            <p:ph type="sldImg"/>
          </p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="3" name="Notes Placeholder 2"/>
          <p:cNvSpPr>
            <a:spLocks noGrp="1"/>
          </p:cNvSpPr>
          <p:nvPr>
            <p:ph type="body" idx="1"/>
          </p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          {paragraphs_xml}
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>
</p:notes>'''


def create_notes_slide_rels_xml(slide_num: int) -> str:
    """
    创建备注幻灯片关系文件 XML
    
    Args:
        slide_num: 幻灯片序号
    
    Returns:
        关系文件 XML 字符串
    """
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_num}.xml"/>
</Relationships>'''


def create_slide_xml_with_svg(
    slide_num: int, 
    png_rid: str,
    svg_rid: str, 
    width_emu: int, 
    height_emu: int,
    transition: Optional[str] = None,
    transition_duration: float = 0.5,
    auto_advance: Optional[float] = None,
    use_compat_mode: bool = True
) -> str:
    """
    创建包含 SVG 图片的幻灯片 XML
    
    Args:
        slide_num: 幻灯片序号
        png_rid: PNG 后备图片关系 ID
        svg_rid: SVG 关系 ID
        width_emu: 宽度（EMU）
        height_emu: 高度（EMU）
        transition: 切换效果名称
        transition_duration: 切换持续时间（秒）
        auto_advance: 自动翻页间隔（秒）
        use_compat_mode: 是否使用兼容模式（PNG + SVG 双格式）
    """
    # 生成切换效果 XML
    transition_xml = ''
    if transition and ANIMATIONS_AVAILABLE:
        transition_xml = '\n' + create_transition_xml(
            effect=transition,
            duration=transition_duration,
            advance_after=auto_advance
        )
    
    # 兼容模式：PNG 主图片 + SVG 扩展（Office 官方推荐）
    if use_compat_mode:
        blip_xml = f'''<a:blip r:embed="{png_rid}">
            <a:extLst>
              <a:ext uri="{{96DAC541-7B7A-43D3-8B79-37D633B846F1}}">
                <asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" r:embed="{svg_rid}"/>
              </a:ext>
            </a:extLst>
          </a:blip>'''
    else:
        # 纯 SVG 模式（仅新版 Office 支持）
        blip_xml = f'<a:blip r:embed="{svg_rid}"/>'
    
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:pic>
        <p:nvPicPr>
          <p:cNvPr id="2" name="SVG Image {slide_num}"/>
          <p:cNvPicPr>
            <a:picLocks noChangeAspect="1"/>
          </p:cNvPicPr>
          <p:nvPr/>
        </p:nvPicPr>
        <p:blipFill>
          {blip_xml}
          <a:stretch>
            <a:fillRect/>
          </a:stretch>
        </p:blipFill>
        <p:spPr>
          <a:xfrm>
            <a:off x="0" y="0"/>
            <a:ext cx="{width_emu}" cy="{height_emu}"/>
          </a:xfrm>
          <a:prstGeom prst="rect">
            <a:avLst/>
          </a:prstGeom>
        </p:spPr>
      </p:pic>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>{transition_xml}
</p:sld>'''


def create_slide_rels_xml(png_rid: str, png_filename: str, svg_rid: str, svg_filename: str, use_compat_mode: bool = True) -> str:
    """
    创建幻灯片关系文件 XML
    
    Args:
        png_rid: PNG 图片关系 ID
        png_filename: PNG 文件名
        svg_rid: SVG 关系 ID  
        svg_filename: SVG 文件名
        use_compat_mode: 是否使用兼容模式
    """
    if use_compat_mode:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="{png_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{png_filename}"/>
  <Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{svg_filename}"/>
</Relationships>'''
    else:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{svg_filename}"/>
</Relationships>'''


def create_pptx_with_native_svg(
    svg_files: List[Path],
    output_path: Path,
    canvas_format: Optional[str] = None,
    verbose: bool = True,
    transition: Optional[str] = None,
    transition_duration: float = 0.5,
    auto_advance: Optional[float] = None,
    use_compat_mode: bool = True,
    notes: Optional[dict] = None,
    enable_notes: bool = True
) -> bool:
    """
    创建包含原生 SVG 的 PPTX 文件
    
    Args:
        svg_files: SVG 文件列表
        output_path: 输出路径
        canvas_format: 画布格式
        verbose: 是否输出详细信息
        transition: 切换效果 (fade/push/wipe/split/reveal/cover/random)
        transition_duration: 切换持续时间（秒）
        auto_advance: 自动翻页间隔（秒）
        use_compat_mode: 使用 Office 兼容模式（PNG + SVG 双格式，默认开启）
        notes: 备注字典，key 为幻灯片编号，value 为备注内容
        enable_notes: 是否启用备注嵌入（默认开启）
    """
    if not svg_files:
        print("错误: 没有找到 SVG 文件")
        return False
    
    # 检查兼容模式依赖
    renderer_name, renderer_status, renderer_hint = get_png_renderer_info()
    if use_compat_mode and PNG_RENDERER is None:
        print("警告: 未安装 PNG 渲染库，无法使用兼容模式")
        print(f"  {renderer_hint}")
        print("  将使用纯 SVG 模式（可能在 Office LTSC 2021 等版本中不显示）")
        use_compat_mode = False
    
    # 自动检测画布格式或从 viewBox 获取尺寸
    custom_pixels: Optional[Tuple[int, int]] = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  检测到画布格式: {format_name}")
    
    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  使用 SVG viewBox 尺寸: {custom_pixels[0]} x {custom_pixels[1]} px")
    
    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  使用默认格式: PPT 16:9")
    
    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)
    
    if verbose:
        print(f"  幻灯片尺寸: {pixel_width} x {pixel_height} px")
        print(f"  SVG 文件数: {len(svg_files)}")
        if use_compat_mode:
            print(f"  兼容模式: 开启 (PNG + SVG 双格式)")
            print(f"  PNG 渲染: {renderer_name} {renderer_status}")
        else:
            print(f"  兼容模式: 关闭 (纯 SVG)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  切换效果: {trans_name}")
        if enable_notes and notes:
            print(f"  演讲备注: {len(notes)} 页")
        elif enable_notes:
            print(f"  演讲备注: 已启用（未找到备注文件）")
        else:
            print(f"  演讲备注: 已禁用")
        print()
    
    # 创建临时目录
    temp_dir = Path(tempfile.mkdtemp())
    
    try:
        # 首先用 python-pptx 创建基础 PPTX
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
        
        # 添加空白幻灯片作为占位
        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)
        
        # 保存基础 PPTX
        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))
        
        # 解压 PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)
        
        # 创建 media 目录
        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)
        
        # 处理每个 SVG 文件
        success_count = 0
        any_png_generated = False
        
        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i
            svg_filename = f'image{i}.svg'
            png_filename = f'image{i}.png'
            png_rid = 'rId2'
            svg_rid = 'rId3' if use_compat_mode else 'rId2'
            
            try:
                # 复制 SVG 到 media 目录
                shutil.copy(svg_path, media_dir / svg_filename)
                
                # 兼容模式：生成 PNG 后备图片
                slide_has_png = False
                if use_compat_mode:
                    png_path = media_dir / png_filename
                    png_success = convert_svg_to_png(
                        svg_path, 
                        png_path,
                        width=pixel_width,
                        height=pixel_height
                    )
                    if png_success:
                        slide_has_png = True
                        any_png_generated = True
                    else:
                        # PNG 生成失败，降级为纯 SVG
                        if verbose:
                            print(f"  [{i}/{len(svg_files)}] {svg_path.name} - PNG 生成失败，使用纯 SVG")
                        svg_rid = 'rId2'
                
                # 更新幻灯片 XML
                slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                slide_xml = create_slide_xml_with_svg(
                    slide_num, 
                    png_rid=png_rid,
                    svg_rid=svg_rid, 
                    width_emu=width_emu, 
                    height_emu=height_emu,
                    transition=transition,
                    transition_duration=transition_duration,
                    auto_advance=auto_advance,
                    use_compat_mode=(use_compat_mode and slide_has_png)
                )
                with open(slide_xml_path, 'w', encoding='utf-8') as f:
                    f.write(slide_xml)
                
                # 创建/更新关系文件
                rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                rels_dir.mkdir(exist_ok=True)
                rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                rels_xml = create_slide_rels_xml(
                    png_rid=png_rid,
                    png_filename=png_filename,
                    svg_rid=svg_rid, 
                    svg_filename=svg_filename,
                    use_compat_mode=(use_compat_mode and slide_has_png)
                )
                with open(rels_path, 'w', encoding='utf-8') as f:
                    f.write(rels_xml)
                
                # 处理备注
                notes_content = ''
                if enable_notes:
                    # 按文件名匹配（新逻辑）或按序号匹配（向后兼容）
                    svg_stem = svg_path.stem
                    notes_content = notes.get(svg_stem, '') if notes else ''
                    if notes_content:
                        notes_text = markdown_to_plain_text(notes_content)
                    else:
                        notes_text = ''  # 空备注
                    
                    # 创建 notesSlides 目录
                    notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                    notes_slides_dir.mkdir(exist_ok=True)
                    
                    # 创建备注幻灯片 XML
                    notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                    notes_xml = create_notes_slide_xml(slide_num, notes_text)
                    with open(notes_xml_path, 'w', encoding='utf-8') as f:
                        f.write(notes_xml)
                    
                    # 创建备注幻灯片关系文件
                    notes_rels_dir = notes_slides_dir / '_rels'
                    notes_rels_dir.mkdir(exist_ok=True)
                    notes_rels_path = notes_rels_dir / f'notesSlide{slide_num}.xml.rels'
                    notes_rels_xml = create_notes_slide_rels_xml(slide_num)
                    with open(notes_rels_path, 'w', encoding='utf-8') as f:
                        f.write(notes_rels_xml)
                    
                    # 更新 slide.xml.rels 添加备注关联
                    with open(rels_path, 'r', encoding='utf-8') as f:
                        slide_rels_content = f.read()
                    notes_rel = f'  <Relationship Id="rId10" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{slide_num}.xml"/>'
                    slide_rels_content = slide_rels_content.replace('</Relationships>', notes_rel + '\n</Relationships>')
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(slide_rels_content)
                
                if verbose:
                    mode_str = " (PNG+SVG)" if (use_compat_mode and slide_has_png) else " (SVG)"
                    has_notes = enable_notes and bool(notes_content)
                    notes_str = " +备注" if has_notes else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}")
                
                success_count += 1
                
            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - 错误: {e}")
        
        # 更新 [Content_Types].xml 添加 SVG 和 PNG 类型
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()
        
        # 添加 SVG 扩展类型（如果不存在）
        types_to_add = []
        if 'Extension="svg"' not in content_types:
            types_to_add.append('  <Default Extension="svg" ContentType="image/svg+xml"/>')
        if any_png_generated and 'Extension="png"' not in content_types:
            types_to_add.append('  <Default Extension="png" ContentType="image/png"/>')
        
        if types_to_add:
            content_types = content_types.replace(
                '</Types>',
                '\n'.join(types_to_add) + '\n</Types>'
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)
        
        # 添加 notesSlides 内容类型
        if enable_notes:
            for i in range(1, len(svg_files) + 1):
                override = f'  <Override PartName="/ppt/notesSlides/notesSlide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                if override not in content_types:
                    content_types = content_types.replace('</Types>', override + '\n</Types>')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)
        
        # 重新打包 PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)
        
        if verbose:
            print()
            print(f"[完成] 已保存: {output_path}")
            print(f"  成功: {success_count}, 失败: {len(svg_files) - success_count}")
            if use_compat_mode and any_png_generated:
                print(f"  模式: Office 兼容模式 (支持所有 Office 版本)")
                # 如果使用 svglib，给出升级提示
                if PNG_RENDERER == 'svglib' and renderer_hint:
                    print(f"  [提示] {renderer_hint}")
        
        return success_count == len(svg_files)
        
    finally:
        # 清理临时目录
        shutil.rmtree(temp_dir, ignore_errors=True)


def main():
    # 构建切换效果选项列表
    transition_choices = list(TRANSITIONS.keys()) if TRANSITIONS else ['fade', 'push', 'wipe', 'split', 'reveal', 'cover', 'random']
    
    parser = argparse.ArgumentParser(
        description='PPT Master - SVG 转 PPTX 工具（Office 兼容模式）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=f'''
示例:
    %(prog)s examples/ppt169_demo -s final    # 推荐：使用后处理完成的版本
    %(prog)s examples/ppt169_demo             # 使用原始版本
    %(prog)s examples/ppt169_demo -o presentation.pptx
    %(prog)s examples/ppt169_demo --no-compat # 禁用兼容模式（仅纯 SVG）
    
    # 添加页面切换效果
    %(prog)s examples/ppt169_demo --transition fade
    %(prog)s examples/ppt169_demo -t push --transition-duration 1.0

SVG 来源目录 (-s):
    output   - svg_output（原始版本）
    final    - svg_final（后处理完成，推荐）
    <任意名> - 直接指定项目下的子目录名

切换效果 (-t/--transition):
    {', '.join(transition_choices)}

兼容模式 (默认开启):
    - 自动生成 PNG 后备图片，SVG 作为扩展嵌入
    - 兼容所有 Office 版本（包括 Office LTSC 2021）
    - 新版 Office 仍显示 SVG（可编辑），旧版显示 PNG
    - 需要安装 svglib: pip install svglib reportlab
    - 使用 --no-compat 可禁用（仅 Office 2019+ 支持）

演讲备注 (默认开启):
    - 自动读取 notes/ 目录中的 Markdown 备注文件
    - 支持两种命名方式：
      1. 按文件名匹配（推荐）：01_封面.md 对应 01_封面.svg
      2. 按序号匹配：slide01.md 对应第1个 SVG（向后兼容）
    - 使用 --no-notes 可禁用
'''
    )
    
    parser.add_argument('project_path', type=str, help='项目目录路径')
    parser.add_argument('-o', '--output', type=str, default=None, help='输出文件路径')
    parser.add_argument('-s', '--source', type=str, default='output', 
                        help='SVG 来源: output/final 或任意子目录名 (推荐 final)')
    parser.add_argument('-f', '--format', type=str, choices=list(CANVAS_FORMATS.keys()), default=None, help='指定画布格式')
    parser.add_argument('-q', '--quiet', action='store_true', help='静默模式')
    
    # 兼容模式参数
    parser.add_argument('--no-compat', action='store_true',
                        help='禁用 Office 兼容模式（仅使用纯 SVG，需要 Office 2019+）')
    
    # 切换效果参数
    parser.add_argument('-t', '--transition', type=str, choices=transition_choices, default=None,
                        help='页面切换效果 (默认: 无)')
    parser.add_argument('--transition-duration', type=float, default=0.5,
                        help='切换持续时间/秒 (默认: 0.5)')
    parser.add_argument('--auto-advance', type=float, default=None,
                        help='自动翻页间隔/秒 (默认: 手动翻页)')
    
    # 备注参数
    parser.add_argument('--no-notes', action='store_true',
                        help='禁用演讲备注嵌入（默认启用）')
    
    args = parser.parse_args()
    
    project_path = Path(args.project_path)
    if not project_path.exists():
        print(f"错误: 路径不存在: {project_path}")
        sys.exit(1)
    
    try:
        project_info = get_project_info(str(project_path))
        project_name = project_info.get('name', project_path.name)
        detected_format = project_info.get('format')
    except Exception:
        project_name = project_path.name
        detected_format = None
    
    canvas_format = args.format
    if canvas_format is None and detected_format and detected_format != 'unknown':
        canvas_format = detected_format
    
    svg_files, source_dir_name = find_svg_files(project_path, args.source)
    
    if not svg_files:
        print("错误: 未找到 SVG 文件")
        sys.exit(1)
    
    if args.output:
        output_path = Path(args.output)
    else:
        # 默认带时间戳，避免覆盖之前的版本
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = project_path / f"{project_name}_{timestamp}.pptx"
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    verbose = not args.quiet
    
    # 读取备注文件
    enable_notes = not args.no_notes
    notes = {}
    if enable_notes:
        notes = find_notes_files(project_path, svg_files)
    
    if verbose:
        print("PPT Master - SVG 转 PPTX 工具（原生 SVG）")
        print("=" * 50)
        print(f"  项目路径: {project_path}")
        print(f"  SVG 目录: {source_dir_name}")
        print(f"  输出文件: {output_path}")
        print()
    
    success = create_pptx_with_native_svg(
        svg_files,
        output_path,
        canvas_format=canvas_format,
        verbose=verbose,
        transition=args.transition,
        transition_duration=args.transition_duration,
        auto_advance=args.auto_advance,
        use_compat_mode=not args.no_compat,
        notes=notes,
        enable_notes=enable_notes
    )
    
    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
